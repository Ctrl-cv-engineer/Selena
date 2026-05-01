"""Document generation skill runtime — produces PDF, DOCX, and PPTX files."""

from __future__ import annotations

import json
import logging
import os
import time

logger = logging.getLogger(__name__)

try:
    from DialogueSystem.config.paths import DATA_DIR
except ImportError:
    try:
        from DialogueSystem.config.paths import DATA_DIR
    except ImportError:
        DATA_DIR = os.path.join(os.path.dirname(os.path.abspath(__file__)), "..", "data")

DOCUMENTS_DIR = os.path.join(DATA_DIR, "documents")


def _ensure_documents_dir() -> str:
    os.makedirs(DOCUMENTS_DIR, exist_ok=True)
    return DOCUMENTS_DIR


def _resolve_output_path(output_path: str, title: str, fmt: str) -> str:
    if output_path:
        resolved = os.path.abspath(output_path)
        parent = os.path.dirname(resolved)
        if parent:
            os.makedirs(parent, exist_ok=True)
        return resolved
    _ensure_documents_dir()
    safe_title = "".join(c if c.isalnum() or c in (" ", "-", "_") else "_" for c in (title or "document"))[:60].strip()
    timestamp = time.strftime("%Y%m%d_%H%M%S")
    return os.path.join(DOCUMENTS_DIR, f"{safe_title}_{timestamp}.{fmt}")


def _parse_content(content) -> dict:
    if isinstance(content, str):
        return json.loads(content)
    return dict(content or {})


def _generate_docx(content: dict, output_path: str) -> str:
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    title = str(content.get("title") or "Untitled").strip()
    doc.add_heading(title, level=0)

    author = str(content.get("author") or "").strip()
    date_str = str(content.get("date") or "").strip()
    if author or date_str:
        meta_parts = []
        if author:
            meta_parts.append(f"Author: {author}")
        if date_str:
            meta_parts.append(f"Date: {date_str}")
        p = doc.add_paragraph(" | ".join(meta_parts))
        p.style.font.size = Pt(10)

    for section in content.get("sections") or []:
        heading = str(section.get("heading") or "").strip()
        body = str(section.get("body") or "").strip()
        level = min(max(int(section.get("level") or 1), 1), 3)
        if heading:
            doc.add_heading(heading, level=level)
        if body:
            for paragraph in body.split("\n"):
                paragraph = paragraph.strip()
                if paragraph:
                    doc.add_paragraph(paragraph)

    doc.save(output_path)
    return output_path


def _generate_pdf(content: dict, output_path: str) -> str:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer

    doc = SimpleDocTemplate(
        output_path,
        pagesize=A4,
        leftMargin=2.5 * cm,
        rightMargin=2.5 * cm,
        topMargin=2 * cm,
        bottomMargin=2 * cm,
    )
    styles = getSampleStyleSheet()
    title_style = ParagraphStyle(
        "DocTitle", parent=styles["Title"], fontSize=20, spaceAfter=12,
    )
    heading_styles = {
        1: ParagraphStyle("H1", parent=styles["Heading1"], fontSize=16, spaceAfter=8),
        2: ParagraphStyle("H2", parent=styles["Heading2"], fontSize=13, spaceAfter=6),
        3: ParagraphStyle("H3", parent=styles["Heading3"], fontSize=11, spaceAfter=4),
    }
    body_style = styles["BodyText"]

    story = []
    title = str(content.get("title") or "Untitled").strip()
    story.append(Paragraph(title, title_style))

    author = str(content.get("author") or "").strip()
    date_str = str(content.get("date") or "").strip()
    if author or date_str:
        meta = " | ".join(p for p in [author, date_str] if p)
        story.append(Paragraph(meta, styles["Normal"]))
    story.append(Spacer(1, 12))

    for section in content.get("sections") or []:
        heading = str(section.get("heading") or "").strip()
        body = str(section.get("body") or "").strip()
        level = min(max(int(section.get("level") or 1), 1), 3)
        if heading:
            story.append(Paragraph(heading, heading_styles[level]))
        if body:
            for paragraph in body.split("\n"):
                paragraph = paragraph.strip()
                if paragraph:
                    story.append(Paragraph(paragraph, body_style))
            story.append(Spacer(1, 6))

    doc.build(story)
    return output_path


def _generate_pptx(content: dict, output_path: str) -> str:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    title = str(content.get("title") or "Untitled").strip()

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = title
    subtitle_placeholder = title_slide.placeholders[1] if len(title_slide.placeholders) > 1 else None
    if subtitle_placeholder:
        meta_parts = []
        author = str(content.get("author") or "").strip()
        date_str = str(content.get("date") or "").strip()
        if author:
            meta_parts.append(author)
        if date_str:
            meta_parts.append(date_str)
        subtitle_placeholder.text = " | ".join(meta_parts) if meta_parts else ""

    for section in content.get("sections") or []:
        heading = str(section.get("heading") or "").strip()
        body = str(section.get("body") or "").strip()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = heading or "Slide"
        body_placeholder = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        if body_placeholder and body:
            tf = body_placeholder.text_frame
            tf.clear()
            for i, line in enumerate(body.split("\n")):
                line = line.strip()
                if not line:
                    continue
                if i == 0:
                    tf.paragraphs[0].text = line
                    tf.paragraphs[0].font.size = Pt(18)
                else:
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(18)

    prs.save(output_path)
    return output_path


GENERATORS = {
    "docx": _generate_docx,
    "pdf": _generate_pdf,
    "pptx": _generate_pptx,
}


def generate_document(Format: str, Content, OutputPath: str = ""):
    normalized_format = str(Format or "").strip().lower()
    if normalized_format not in GENERATORS:
        return {"ok": False, "error": f"Unsupported format: {Format}. Use pdf, docx, or pptx."}

    try:
        parsed_content = _parse_content(Content)
    except Exception as error:
        return {"ok": False, "error": f"Invalid Content: {error}"}

    title = str(parsed_content.get("title") or "document").strip()
    output_path = _resolve_output_path(OutputPath, title, normalized_format)

    try:
        result_path = GENERATORS[normalized_format](parsed_content, output_path)
        file_size = os.path.getsize(result_path)
        return {
            "ok": True,
            "format": normalized_format,
            "output_path": result_path,
            "file_size": file_size,
            "title": title,
        }
    except Exception as error:
        logger.exception("Document generation failed | format=%s", normalized_format)
        return {"ok": False, "error": f"Generation failed: {error}"}


def register_tools(registry):
    registry.register("generateDocument", generate_document)
